#!/usr/bin/env python3
"""
build_pptx.py — Gerador de apresentações Entrelaços Psicologia

Recebe um JSON estruturado com slides e gera um arquivo .pptx 16:9 (1920x1080) aplicando
a identidade visual Entrelaços (paleta, tipografia, layouts, elementos decorativos).

Uso:
    python build_pptx.py --input slides.json --output minha-aula.pptx [--theme dark|light|warm]

Schema do JSON: ver scripts/slide_schema.json

Dependências:
    pip install python-pptx --break-system-packages
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("ERRO: python-pptx não instalado.")
    print("Instale com: pip install python-pptx --break-system-packages")
    sys.exit(1)


# ============================================================
# DESIGN TOKENS — Identidade Entrelaços
# ============================================================

THEMES = {
    "dark": {
        "bg_primary": RGBColor(0x0A, 0x0A, 0x0F),
        "bg_secondary": RGBColor(0x0F, 0x0F, 0x1A),
        "bg_card": RGBColor(0x16, 0x16, 0x25),
        "bg_card_light": RGBColor(0x1E, 0x1E, 0x35),
        "purple": RGBColor(0x6B, 0x46, 0xC1),
        "purple_light": RGBColor(0x8B, 0x5C, 0xF6),
        "purple_dark": RGBColor(0x4C, 0x1D, 0x95),
        "purple_glow": RGBColor(0x7C, 0x3A, 0xED),
        "orange": RGBColor(0xF9, 0x73, 0x16),
        "orange_light": RGBColor(0xFB, 0x92, 0x3C),
        "teal": RGBColor(0x14, 0xB8, 0xA6),
        "text_primary": RGBColor(0xFF, 0xFF, 0xFF),
        "text_secondary": RGBColor(0xCB, 0xD5, 0xE1),
        "text_muted": RGBColor(0x64, 0x74, 0x8B),
        "border": RGBColor(0x2D, 0x2D, 0x50),
    },
    "light": {
        "bg_primary": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_secondary": RGBColor(0xF8, 0xFA, 0xFC),
        "bg_card": RGBColor(0xED, 0xE9, 0xFE),
        "bg_card_light": RGBColor(0xF8, 0xFA, 0xFC),
        "purple": RGBColor(0x6B, 0x46, 0xC1),
        "purple_light": RGBColor(0x8B, 0x5C, 0xF6),
        "purple_dark": RGBColor(0x4C, 0x1D, 0x95),
        "purple_glow": RGBColor(0x7C, 0x3A, 0xED),
        "orange": RGBColor(0xF9, 0x73, 0x16),
        "orange_light": RGBColor(0xFB, 0x92, 0x3C),
        "teal": RGBColor(0x14, 0xB8, 0xA6),
        "text_primary": RGBColor(0x0F, 0x17, 0x2A),
        "text_secondary": RGBColor(0x33, 0x41, 0x55),
        "text_muted": RGBColor(0x64, 0x74, 0x8B),
        "border": RGBColor(0xE2, 0xE8, 0xF0),
    },
    "warm": {
        "bg_primary": RGBColor(0xFA, 0xF7, 0xF2),
        "bg_secondary": RGBColor(0xF5, 0xED, 0xE0),
        "bg_card": RGBColor(0xE7, 0xE0, 0xD2),
        "bg_card_light": RGBColor(0xF0, 0xE9, 0xDC),
        "purple": RGBColor(0x6B, 0x46, 0xC1),
        "purple_light": RGBColor(0x8B, 0x5C, 0xF6),
        "purple_dark": RGBColor(0x4C, 0x1D, 0x95),
        "purple_glow": RGBColor(0x7C, 0x3A, 0xED),
        "orange": RGBColor(0xB8, 0x50, 0x42),  # terracota
        "orange_light": RGBColor(0xC9, 0x6E, 0x60),
        "teal": RGBColor(0xA7, 0xBE, 0xAE),  # sage
        "text_primary": RGBColor(0x2C, 0x18, 0x10),
        "text_secondary": RGBColor(0x5A, 0x3D, 0x2E),
        "text_muted": RGBColor(0x8B, 0x6B, 0x52),
        "border": RGBColor(0xD9, 0xCC, 0xB8),
    },
}

# Fontes (instaladas como Google Fonts → assumir que estão no sistema; fallback Calibri)
FONT_DISPLAY = "Syne"
FONT_DISPLAY_FALLBACK = "Arial Black"
FONT_BODY = "DM Sans"
FONT_BODY_FALLBACK = "Calibri"

# Slide 16:9 — 1920×1080 = 13.333" × 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.667)  # ~96px

# ============================================================
# HELPERS DE BAIXO NÍVEL
# ============================================================


def _set_bg(slide, color):
    """Define cor sólida de fundo do slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, fill_color, line_color=None):
    """Adiciona retângulo sólido."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _add_rounded_rect(slide, x, y, w, h, fill_color, radius=Inches(0.15)):
    """Retângulo com cantos arredondados."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_circle(slide, cx, cy, diameter, fill_color):
    """Círculo centrado em (cx, cy)."""
    r = diameter / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_text(slide, x, y, w, h, text, *,
              font=FONT_BODY, size=18, bold=False, italic=False,
              color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.3, letter_spacing=None):
    """Adiciona caixa de texto com formatação completa."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color

    # letter-spacing via XML (spc em centésimos de ponto)
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(letter_spacing * 100)))

    return tb


def _add_butterfly(slide, x, y, size, theme):
    """Borboleta decorativa estilizada (versão simplificada via shapes)."""
    # Asas — duas elipses sobrepostas em roxo
    wing_color = theme["purple"]
    body_color = theme["orange"]

    # Asa esquerda
    left = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + size * 0.1,
                                   size * 0.45, size * 0.55)
    left.fill.solid()
    left.fill.fore_color.rgb = wing_color
    left.line.fill.background()
    left.rotation = -20

    # Asa direita
    right = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + size * 0.55, y + size * 0.1,
                                    size * 0.45, size * 0.55)
    right.fill.solid()
    right.fill.fore_color.rgb = wing_color
    right.line.fill.background()
    right.rotation = 20

    # Corpo central
    body = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + size * 0.42, y + size * 0.25,
                                   size * 0.16, size * 0.4)
    body.fill.solid()
    body.fill.fore_color.rgb = body_color
    body.line.fill.background()

    # Aplicar transparência via XML
    for shape in [left, right]:
        _set_shape_opacity(shape, 25)  # 25% opacity


def _set_shape_opacity(shape, percent):
    """Define opacidade de uma shape via XML."""
    sp = shape.fill._xPr
    fill_elem = sp.find(qn('a:solidFill'))
    if fill_elem is not None:
        rgb = fill_elem.find(qn('a:srgbClr'))
        if rgb is not None:
            alpha = etree.SubElement(rgb, qn('a:alpha'))
            alpha.set('val', str(int(percent * 1000)))


def _add_section_tag(slide, number, label, theme, y=Inches(0.5)):
    """Tag de seção '/// 02 LABEL'."""
    text = f"///   {number:02d}   {label.upper()}"
    _add_text(slide, MARGIN, y, Inches(8), Inches(0.4),
              text, font=FONT_BODY, size=12, bold=True,
              color=theme["purple_light"], letter_spacing=2.5)


def _add_footer(slide, page_num, total, theme):
    """Footer minimalista nos slides intermediários."""
    _add_text(slide, MARGIN, SLIDE_H - Inches(0.45),
              Inches(6), Inches(0.3),
              "Entrelaços Psicologia",
              font=FONT_BODY, size=10, color=theme["text_muted"])
    _add_text(slide, SLIDE_W - MARGIN - Inches(2), SLIDE_H - Inches(0.45),
              Inches(2), Inches(0.3),
              f"{page_num:02d} / {total:02d}",
              font=FONT_BODY, size=10, color=theme["text_muted"],
              align=PP_ALIGN.RIGHT)


# ============================================================
# RENDERERS POR LAYOUT
# ============================================================


def render_cover(slide, data, theme):
    _set_bg(slide, theme["bg_primary"])

    # Glow radial roxo (simulado via círculo grande translúcido)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   SLIDE_W * 0.6, -Inches(2),
                                   Inches(8), Inches(8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = theme["purple"]
    glow.line.fill.background()
    _set_shape_opacity(glow, 15)

    # Tag (opcional)
    if data.get("tag"):
        _add_text(slide, MARGIN, MARGIN,
                  Inches(8), Inches(0.4),
                  data["tag"].upper(),
                  font=FONT_BODY, size=14, bold=True,
                  color=theme["orange"], letter_spacing=2.5)

    # Título
    title_y = Inches(2.4)
    _add_text(slide, MARGIN, title_y, Inches(10.5), Inches(2.5),
              data["title"],
              font=FONT_DISPLAY, size=72, bold=True,
              color=theme["text_primary"], line_spacing=1.05)

    # Subtítulo
    if data.get("subtitle"):
        _add_text(slide, MARGIN, title_y + Inches(2.3),
                  Inches(10.5), Inches(1.2),
                  data["subtitle"],
                  font=FONT_BODY, size=24,
                  color=theme["text_secondary"], line_spacing=1.4)

    # Autor / assinatura no canto inferior esquerdo
    author = data.get("author", "Tati Ribeiro · Entrelaços Psicologia")
    _add_text(slide, MARGIN, SLIDE_H - Inches(0.85),
              Inches(8), Inches(0.4),
              author,
              font=FONT_BODY, size=14, bold=True,
              color=theme["purple_light"], letter_spacing=1.5)

    # Borboleta decorativa
    _add_butterfly(slide, SLIDE_W - Inches(2.8), Inches(0.6),
                   Inches(2.2), theme)


def render_agenda(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    if data.get("section"):
        _add_section_tag(slide, data.get("section_number", 0),
                         data["section"], theme)

    # Título
    _add_text(slide, MARGIN, Inches(1.2), Inches(11), Inches(1.2),
              data["title"],
              font=FONT_DISPLAY, size=48, bold=True,
              color=theme["text_primary"], line_spacing=1.1)

    # Lista numerada
    items = data.get("items", [])
    item_h = Inches(0.7)
    start_y = Inches(2.8)

    for i, item in enumerate(items):
        y = start_y + (item_h + Inches(0.15)) * i

        # Círculo numerado
        circle_size = Inches(0.55)
        cx = MARGIN + circle_size / 2
        cy = y + circle_size / 2
        circ = _add_circle(slide, cx, cy, circle_size, theme["purple"])

        _add_text(slide, MARGIN, y, circle_size, circle_size,
                  str(i + 1).zfill(2),
                  font=FONT_DISPLAY, size=18, bold=True,
                  color=theme["text_primary"],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Texto do item
        _add_text(slide, MARGIN + Inches(1.0), y + Inches(0.05),
                  Inches(10), Inches(0.6),
                  item,
                  font=FONT_BODY, size=22,
                  color=theme["text_primary"],
                  anchor=MSO_ANCHOR.MIDDLE)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_concept(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    if data.get("section"):
        _add_section_tag(slide, data.get("section_number", 0),
                         data["section"], theme)

    # Título grande
    _add_text(slide, MARGIN, Inches(1.5), Inches(11), Inches(2.2),
              data["title"],
              font=FONT_DISPLAY, size=46, bold=True,
              color=theme["text_primary"], line_spacing=1.15)

    # Body abaixo
    if data.get("body"):
        _add_text(slide, MARGIN, Inches(4.2), Inches(10.5), Inches(2.5),
                  data["body"],
                  font=FONT_BODY, size=22,
                  color=theme["text_secondary"], line_spacing=1.5)

    # Linha decorativa lateral roxa
    _add_rect(slide, MARGIN, Inches(1.5), Inches(0.08), Inches(2.0),
              theme["orange"])

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_quote(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"] if theme is THEMES["dark"]
            else theme["bg_secondary"])

    # Aspas decorativas grandes
    _add_text(slide, MARGIN, Inches(0.8), Inches(2.5), Inches(2.5),
              "“",  # aspas tipográficas
              font=FONT_DISPLAY, size=180, bold=True,
              color=theme["purple"])

    # Citação
    _add_text(slide, Inches(2), Inches(2.3), Inches(9.5), Inches(3.5),
              data["quote"],
              font=FONT_DISPLAY, size=38, bold=True, italic=True,
              color=theme["text_primary"], line_spacing=1.3,
              align=PP_ALIGN.LEFT)

    # Atribuição
    if data.get("attribution"):
        _add_text(slide, Inches(2), SLIDE_H - Inches(1.6),
                  Inches(9.5), Inches(0.6),
                  f"— {data['attribution']}",
                  font=FONT_BODY, size=18,
                  color=theme["purple_light"])

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_two_column(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    # Título no topo
    _add_text(slide, MARGIN, Inches(0.7), Inches(11), Inches(1),
              data.get("title", ""),
              font=FONT_DISPLAY, size=36, bold=True,
              color=theme["text_primary"])

    col_w = (SLIDE_W - MARGIN * 2 - Inches(0.5)) / 2
    col_y = Inches(2.2)
    col_h = Inches(4.5)

    # Coluna esquerda
    left = data.get("left", {})
    _add_rounded_rect(slide, MARGIN, col_y, col_w, col_h, theme["bg_card"])
    _add_text(slide, MARGIN + Inches(0.4), col_y + Inches(0.4),
              col_w - Inches(0.8), Inches(0.6),
              left.get("title", ""),
              font=FONT_DISPLAY, size=22, bold=True,
              color=theme["orange"])
    items = left.get("items", [])
    for i, it in enumerate(items):
        _add_text(slide,
                  MARGIN + Inches(0.4),
                  col_y + Inches(1.2) + Inches(0.6) * i,
                  col_w - Inches(0.8), Inches(0.5),
                  f"•  {it}",
                  font=FONT_BODY, size=18,
                  color=theme["text_primary"], line_spacing=1.4)

    # Coluna direita
    right = data.get("right", {})
    rx = MARGIN + col_w + Inches(0.5)
    _add_rounded_rect(slide, rx, col_y, col_w, col_h, theme["bg_card_light"])
    _add_text(slide, rx + Inches(0.4), col_y + Inches(0.4),
              col_w - Inches(0.8), Inches(0.6),
              right.get("title", ""),
              font=FONT_DISPLAY, size=22, bold=True,
              color=theme["purple_light"])
    items = right.get("items", [])
    for i, it in enumerate(items):
        _add_text(slide,
                  rx + Inches(0.4),
                  col_y + Inches(1.2) + Inches(0.6) * i,
                  col_w - Inches(0.8), Inches(0.5),
                  f"•  {it}",
                  font=FONT_BODY, size=18,
                  color=theme["text_primary"], line_spacing=1.4)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_three_grid(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    _add_text(slide, MARGIN, Inches(0.7), Inches(11), Inches(1),
              data.get("title", ""),
              font=FONT_DISPLAY, size=36, bold=True,
              color=theme["text_primary"])

    items = data.get("items", [])[:3]
    gap = Inches(0.4)
    card_w = (SLIDE_W - MARGIN * 2 - gap * 2) / 3
    card_h = Inches(4.3)
    card_y = Inches(2.2)

    accent_colors = [theme["purple"], theme["orange"], theme["teal"]]

    for i, item in enumerate(items):
        cx = MARGIN + (card_w + gap) * i
        _add_rounded_rect(slide, cx, card_y, card_w, card_h, theme["bg_card"])

        # Ícone em circle (decorativo)
        accent = accent_colors[i % 3]
        _add_circle(slide, cx + Inches(0.7), card_y + Inches(0.7),
                    Inches(0.6), accent)

        # Sub-título
        _add_text(slide, cx + Inches(0.4), card_y + Inches(1.5),
                  card_w - Inches(0.8), Inches(0.7),
                  item.get("title", ""),
                  font=FONT_DISPLAY, size=22, bold=True,
                  color=theme["text_primary"], line_spacing=1.2)

        # Descrição
        _add_text(slide, cx + Inches(0.4), card_y + Inches(2.4),
                  card_w - Inches(0.8), Inches(1.7),
                  item.get("body", ""),
                  font=FONT_BODY, size=15,
                  color=theme["text_secondary"], line_spacing=1.4)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_four_grid(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    _add_text(slide, MARGIN, Inches(0.6), Inches(11), Inches(0.9),
              data.get("title", ""),
              font=FONT_DISPLAY, size=34, bold=True,
              color=theme["text_primary"])

    items = data.get("items", [])[:4]
    gap = Inches(0.3)
    card_w = (SLIDE_W - MARGIN * 2 - gap) / 2
    card_h = (Inches(5.5) - gap) / 2
    grid_y = Inches(1.9)

    accent_colors = [theme["purple"], theme["orange"],
                     theme["teal"], theme["purple_light"]]

    for i, item in enumerate(items):
        col, row = i % 2, i // 2
        cx = MARGIN + (card_w + gap) * col
        cy = grid_y + (card_h + gap) * row

        _add_rounded_rect(slide, cx, cy, card_w, card_h, theme["bg_card"])

        # Linha lateral colorida
        _add_rect(slide, cx, cy, Inches(0.08), card_h, accent_colors[i])

        _add_text(slide, cx + Inches(0.4), cy + Inches(0.3),
                  card_w - Inches(0.7), Inches(0.7),
                  item.get("title", ""),
                  font=FONT_DISPLAY, size=20, bold=True,
                  color=theme["text_primary"])
        _add_text(slide, cx + Inches(0.4), cy + Inches(1.1),
                  card_w - Inches(0.7), card_h - Inches(1.4),
                  item.get("body", ""),
                  font=FONT_BODY, size=14,
                  color=theme["text_secondary"], line_spacing=1.4)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_bullets(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    if data.get("section"):
        _add_section_tag(slide, data.get("section_number", 0),
                         data["section"], theme)

    # Título com altura maior pra acomodar quebra em 2 linhas
    _add_text(slide, MARGIN, Inches(1.1), Inches(11.5), Inches(2.0),
              data.get("title", ""),
              font=FONT_DISPLAY, size=38, bold=True,
              color=theme["text_primary"], line_spacing=1.15)

    items = data.get("items", [])
    # Posição Y inicial dos bullets — começa em 3.6 dando respiração ao título
    bullets_start = Inches(3.5)
    bullet_gap = Inches(0.65)

    for i, it in enumerate(items):
        y = bullets_start + bullet_gap * i
        # Bullet roxo decorativo (centralizado na altura média do texto)
        _add_circle(slide, MARGIN + Inches(0.15), y + Inches(0.28),
                    Inches(0.18), theme["purple_light"])
        _add_text(slide, MARGIN + Inches(0.55), y,
                  Inches(11), Inches(0.6),
                  it,
                  font=FONT_BODY, size=20,
                  color=theme["text_primary"], line_spacing=1.3,
                  anchor=MSO_ANCHOR.MIDDLE)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_big_stat(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    # Número gigante centralizado
    stat = data.get("stat", "")
    _add_text(slide, MARGIN, Inches(1.5), SLIDE_W - MARGIN * 2, Inches(3.5),
              stat,
              font=FONT_DISPLAY, size=180, bold=True,
              color=theme["purple"], align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # Label principal
    if data.get("label"):
        _add_text(slide, MARGIN, Inches(5.0),
                  SLIDE_W - MARGIN * 2, Inches(0.7),
                  data["label"],
                  font=FONT_DISPLAY, size=26, bold=True,
                  color=theme["text_primary"], align=PP_ALIGN.CENTER)

    # Sub-linha
    if data.get("subline"):
        _add_text(slide, MARGIN, Inches(5.7),
                  SLIDE_W - MARGIN * 2, Inches(0.5),
                  data["subline"],
                  font=FONT_BODY, size=18,
                  color=theme["text_secondary"], align=PP_ALIGN.CENTER)

    # Fonte
    if data.get("source"):
        _add_text(slide, MARGIN, Inches(6.5),
                  SLIDE_W - MARGIN * 2, Inches(0.4),
                  f"— {data['source']}",
                  font=FONT_BODY, size=12, italic=True,
                  color=theme["text_muted"], align=PP_ALIGN.CENTER)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_step_flow(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    _add_text(slide, MARGIN, Inches(0.8), Inches(11), Inches(1),
              data.get("title", ""),
              font=FONT_DISPLAY, size=36, bold=True,
              color=theme["text_primary"])

    steps = data.get("steps", [])
    n = max(len(steps), 1)
    flow_w = SLIDE_W - MARGIN * 2
    step_w = flow_w / n
    cy = Inches(3.6)
    circle_d = Inches(1.1)
    label_padding = Inches(0.25)  # padding em cada lado da label

    for i, step in enumerate(steps):
        cx = MARGIN + step_w * (i + 0.5)

        # Linha conectora (atrás)
        if i < n - 1:
            line_x = MARGIN + step_w * (i + 0.5) + circle_d / 2
            line_w = step_w - circle_d
            _add_rect(slide, line_x, cy + circle_d / 2 - Inches(0.02),
                      line_w, Inches(0.04),
                      theme["purple_light"])

        # Círculo
        _add_circle(slide, cx, cy + circle_d / 2, circle_d, theme["purple"])

        # Número
        _add_text(slide, cx - circle_d / 2, cy,
                  circle_d, circle_d,
                  str(i + 1).zfill(2),
                  font=FONT_DISPLAY, size=22, bold=True,
                  color=theme["text_primary"],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Label (com padding lateral para não colar uma na outra)
        label_w = step_w - label_padding * 2
        _add_text(slide, cx - label_w / 2, cy + circle_d + Inches(0.35),
                  label_w, Inches(1.6),
                  step,
                  font=FONT_BODY, size=14, bold=True,
                  color=theme["text_primary"], align=PP_ALIGN.CENTER,
                  line_spacing=1.35)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_exercise(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    # Banner colorido com a pergunta
    banner_y = Inches(2.0)
    banner_h = Inches(3.5)
    _add_rounded_rect(slide, MARGIN, banner_y,
                      SLIDE_W - MARGIN * 2, banner_h,
                      theme["purple"])

    # Tag "EXERCÍCIO" / "REFLEXÃO"
    label = data.get("label", "REFLEXÃO")
    _add_text(slide, MARGIN + Inches(0.6), banner_y + Inches(0.5),
              Inches(8), Inches(0.5),
              label.upper(),
              font=FONT_BODY, size=14, bold=True,
              color=theme["orange_light"], letter_spacing=2.5)

    # Pergunta central
    _add_text(slide, MARGIN + Inches(0.6), banner_y + Inches(1.2),
              SLIDE_W - MARGIN * 2 - Inches(1.2), Inches(2.0),
              data.get("question", ""),
              font=FONT_DISPLAY, size=34, bold=True,
              color=theme["text_primary"], line_spacing=1.3)

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_manifesto(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    # Borboleta gigante decorativa no fundo
    _add_butterfly(slide, Inches(8), Inches(0.5), Inches(5), theme)

    # Frase manifesto centralizada
    _add_text(slide, MARGIN, Inches(2), SLIDE_W - MARGIN * 2, Inches(3.5),
              data.get("text", ""),
              font=FONT_DISPLAY, size=64, bold=True,
              color=theme["text_primary"], line_spacing=1.15,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Assinatura
    if data.get("signature"):
        _add_text(slide, MARGIN, Inches(6.0),
                  Inches(11), Inches(0.5),
                  data["signature"],
                  font=FONT_BODY, size=16, italic=True,
                  color=theme["purple_light"])

    if page_num is not None and total is not None:
        _add_footer(slide, page_num, total, theme)


def render_cta(slide, data, theme, page_num=None, total=None):
    _set_bg(slide, theme["bg_primary"])

    # Tag superior
    _add_text(slide, MARGIN, Inches(0.6), Inches(11), Inches(0.4),
              data.get("tag", "OFERTA").upper(),
              font=FONT_BODY, size=14, bold=True,
              color=theme["orange"], letter_spacing=2.5)

    # Título
    _add_text(slide, MARGIN, Inches(1.2), Inches(11), Inches(1.5),
              data.get("title", ""),
              font=FONT_DISPLAY, size=42, bold=True,
              color=theme["text_primary"], line_spacing=1.1)

    # Subtítulo
    if data.get("subtitle"):
        _add_text(slide, MARGIN, Inches(2.7),
                  Inches(11), Inches(0.7),
                  data["subtitle"],
                  font=FONT_BODY, size=20,
                  color=theme["text_secondary"])

    # Lista do que está incluído
    included = data.get("included", [])
    list_y = Inches(3.6)
    for i, it in enumerate(included):
        y = list_y + Inches(0.5) * i
        _add_text(slide, MARGIN, y, Inches(0.5), Inches(0.4),
                  "✓", font=FONT_BODY, size=20, bold=True,
                  color=theme["teal"])
        _add_text(slide, MARGIN + Inches(0.5), y,
                  Inches(8), Inches(0.4),
                  it, font=FONT_BODY, size=18,
                  color=theme["text_primary"])

    # Preço (à direita)
    if data.get("price"):
        _add_rounded_rect(slide, Inches(9.0), Inches(3.6),
                          Inches(3.5), Inches(2.0),
                          theme["bg_card"])
        _add_text(slide, Inches(9.2), Inches(3.8),
                  Inches(3.1), Inches(0.4),
                  "INVESTIMENTO", font=FONT_BODY, size=11, bold=True,
                  color=theme["text_muted"], letter_spacing=2.0,
                  align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(9.2), Inches(4.2),
                  Inches(3.1), Inches(1.0),
                  data["price"],
                  font=FONT_DISPLAY, size=32, bold=True,
                  color=theme["orange"], align=PP_ALIGN.CENTER)
        if data.get("price_note"):
            _add_text(slide, Inches(9.2), Inches(5.2),
                      Inches(3.1), Inches(0.4),
                      data["price_note"],
                      font=FONT_BODY, size=12,
                      color=theme["text_secondary"],
                      align=PP_ALIGN.CENTER)

    # Botão CTA + URL
    btn_y = Inches(6.2)
    btn_w = Inches(6.5)
    btn = _add_rounded_rect(slide, MARGIN, btn_y, btn_w, Inches(0.85),
                             theme["orange"], radius=Inches(0.42))
    _add_text(slide, MARGIN, btn_y, btn_w, Inches(0.85),
              data.get("cta_label", "Quero fazer parte"),
              font=FONT_BODY, size=20, bold=True,
              color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if data.get("url"):
        _add_text(slide, MARGIN + btn_w + Inches(0.3), btn_y + Inches(0.25),
                  Inches(5), Inches(0.4),
                  data["url"],
                  font=FONT_BODY, size=14,
                  color=theme["text_secondary"])


def render_closing(slide, data, theme):
    _set_bg(slide, theme["bg_primary"])

    # Glow superior (atrás da borboleta) — pequeno e contido
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   SLIDE_W / 2 - Inches(2.5),
                                   Inches(0.0),
                                   Inches(5), Inches(5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = theme["purple"]
    glow.line.fill.background()
    _set_shape_opacity(glow, 10)

    # Borboleta menor e mais ao topo (não sobrepõe o texto)
    _add_butterfly(slide, SLIDE_W / 2 - Inches(0.9), Inches(0.6),
                   Inches(1.8), theme)

    # Frase principal — centralizada verticalmente na metade inferior
    _add_text(slide, MARGIN, Inches(3.5),
              SLIDE_W - MARGIN * 2, Inches(1.5),
              data.get("text", "Obrigada."),
              font=FONT_DISPLAY, size=60, bold=True,
              color=theme["text_primary"],
              align=PP_ALIGN.CENTER, line_spacing=1.1,
              anchor=MSO_ANCHOR.MIDDLE)

    # Sub-mensagem
    if data.get("subtext"):
        _add_text(slide, MARGIN, Inches(5.2),
                  SLIDE_W - MARGIN * 2, Inches(0.9),
                  data["subtext"],
                  font=FONT_BODY, size=20,
                  color=theme["text_secondary"],
                  align=PP_ALIGN.CENTER, line_spacing=1.4)

    # Assinatura
    sig = data.get("signature", "Tati Ribeiro · Entrelaços Psicologia")
    _add_text(slide, MARGIN, SLIDE_H - Inches(0.95),
              SLIDE_W - MARGIN * 2, Inches(0.4),
              sig,
              font=FONT_BODY, size=14, bold=True,
              color=theme["purple_light"],
              align=PP_ALIGN.CENTER, letter_spacing=1.5)

    # Redes sociais (opcional)
    if data.get("socials"):
        _add_text(slide, MARGIN, SLIDE_H - Inches(0.5),
                  SLIDE_W - MARGIN * 2, Inches(0.3),
                  data["socials"],
                  font=FONT_BODY, size=12,
                  color=theme["text_muted"],
                  align=PP_ALIGN.CENTER)


# ============================================================
# DISPATCHER
# ============================================================


LAYOUT_RENDERERS = {
    "cover": render_cover,
    "agenda": render_agenda,
    "concept": render_concept,
    "quote": render_quote,
    "two-column": render_two_column,
    "two_column": render_two_column,
    "three-grid": render_three_grid,
    "three_grid": render_three_grid,
    "four-grid": render_four_grid,
    "four_grid": render_four_grid,
    "bullets": render_bullets,
    "big-stat": render_big_stat,
    "big_stat": render_big_stat,
    "step-flow": render_step_flow,
    "step_flow": render_step_flow,
    "exercise": render_exercise,
    "manifesto": render_manifesto,
    "cta": render_cta,
    "closing": render_closing,
}


def build_pptx(deck_data, output_path, theme_name="dark"):
    """Constrói o pptx a partir do JSON estruturado."""
    if theme_name not in THEMES:
        print(f"⚠️ Tema '{theme_name}' desconhecido, usando 'dark'.")
        theme_name = "dark"
    theme = THEMES[theme_name]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # layout em branco

    slides_data = deck_data.get("slides", [])
    total = len(slides_data)

    for idx, sd in enumerate(slides_data):
        layout = sd.get("layout", "concept").lower().replace("_", "-")
        renderer = LAYOUT_RENDERERS.get(
            layout.replace("-", "_"),
            LAYOUT_RENDERERS.get(layout, render_concept)
        )

        slide = prs.slides.add_slide(blank)

        # Renderiza
        try:
            if layout in ("cover", "closing"):
                renderer(slide, sd, theme)
            else:
                renderer(slide, sd, theme, page_num=idx + 1, total=total)
        except Exception as e:
            print(f"⚠️ Erro renderizando slide {idx + 1} (layout={layout}): {e}")
            # Fallback: apenas texto cru
            _set_bg(slide, theme["bg_primary"])
            _add_text(slide, MARGIN, Inches(2),
                      SLIDE_W - MARGIN * 2, Inches(3),
                      f"[Erro: {e}]\n\n{json.dumps(sd, ensure_ascii=False, indent=2)[:500]}",
                      font=FONT_BODY, size=14,
                      color=theme["text_primary"])

    prs.save(output_path)
    print(f"✅ Apresentação salva em: {output_path}")
    print(f"   {total} slides · tema: {theme_name}")


# ============================================================
# CLI
# ============================================================


def main():
    parser = argparse.ArgumentParser(description="Gera PPTX no padrão Entrelaços")
    parser.add_argument("--input", "-i", required=True,
                        help="Caminho do JSON com a estrutura dos slides")
    parser.add_argument("--output", "-o", required=True,
                        help="Caminho do .pptx de saída")
    parser.add_argument("--theme", "-t", default=None,
                        choices=["dark", "light", "warm"],
                        help="Tema visual (override do meta.theme do JSON)")
    args = parser.p